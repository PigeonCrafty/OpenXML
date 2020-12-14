import pptx
from pptx import Presentation


my_file = "c:\\Users\\pigeonz.CZ\\OneDrive - RWS\\AI\\openXML\\_Samples\\AmazonTRA\\zhCN\\03_DTP\\Migrations to AWS - Technical 2.0 Presentation Deck.pptx"

text_runs = []

pre = Presentation(my_file)

slides = pre.slides

for slide in slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)

for i in text_runs:
    print(i)